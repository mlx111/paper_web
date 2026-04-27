from __future__ import annotations

import zipfile
from pathlib import Path
from typing import Any


_CONTENT_TYPES_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slides/slide1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
</Types>
"""

_ROOT_RELS_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>
"""

_PRESENTATION_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst>
    <p:sldMasterId id="2147483648" r:id="rId1"/>
  </p:sldMasterIdLst>
  <p:sldIdLst>
    <p:sldId id="256" r:id="rId2"/>
  </p:sldIdLst>
  <p:sldSz cx="12192000" cy="6858000"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>
"""

_PRESENTATION_RELS_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide1.xml"/>
</Relationships>
"""

_SLIDE_MASTER_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld name="Default Master"><p:spTree/></p:cSld>
  <p:sldLayoutIdLst><p:sldLayoutId id="1" r:id="rId1"/></p:sldLayoutIdLst>
  <p:txStyles/>
</p:sldMaster>
"""

_SLIDE_MASTER_RELS_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>
"""

_SLIDE_LAYOUT_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="title">
  <p:cSld name="Layout"><p:spTree/></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>
"""

_THEME_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Simple Theme">
  <a:themeElements>
    <a:clrScheme name="Simple">
      <a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>
      <a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>
      <a:dk2><a:srgbClr val="1F1F1F"/></a:dk2>
      <a:lt2><a:srgbClr val="F5F5F5"/></a:lt2>
      <a:accent1><a:srgbClr val="2F6BFF"/></a:accent1>
      <a:accent2><a:srgbClr val="0F9D58"/></a:accent2>
      <a:accent3><a:srgbClr val="F4B400"/></a:accent3>
      <a:accent4><a:srgbClr val="DB4437"/></a:accent4>
      <a:accent5><a:srgbClr val="5F6368"/></a:accent5>
      <a:accent6><a:srgbClr val="4285F4"/></a:accent6>
      <a:hlink><a:srgbClr val="0000FF"/></a:hlink>
      <a:folHlink><a:srgbClr val="800080"/></a:folHlink>
    </a:clrScheme>
    <a:fontScheme name="Simple">
      <a:majorFont><a:latin typeface="Arial"/></a:majorFont>
      <a:minorFont><a:latin typeface="Arial"/></a:minorFont>
    </a:fontScheme>
    <a:fmtScheme name="Simple"><a:fillStyleLst/><a:lnStyleLst/><a:effectStyleLst/><a:bgFillStyleLst/></a:fmtScheme>
  </a:themeElements>
</a:theme>
"""

_APP_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
 xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>MyPaperWeb</Application>
</Properties>
"""

_CORE_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties"
 xmlns:dc="http://purl.org/dc/elements/1.1/"
 xmlns:dcterms="http://purl.org/dc/terms/"
 xmlns:dcmitype="http://purl.org/dc/dcmitype/"
 xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>{title}</dc:title>
  <dc:creator>MyPaperWeb</dc:creator>
</cp:coreProperties>
"""


def _xml_escape(text: str) -> str:
    return (
        str(text)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


class PresentationExportService:
    """导出基础版 pptx。"""

    def export(self, plan: dict[str, Any], manuscript: str, output_path: Path) -> Path:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            from pptx import Presentation  # type: ignore

            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = str(plan.get("title", "Presentation"))
            if len(title_slide.placeholders) > 1:
                title_slide.placeholders[1].text = str(plan.get("audience", ""))

            for slide in plan.get("slides", []):
                ppt_slide = prs.slides.add_slide(prs.slide_layouts[1])
                ppt_slide.shapes.title.text = str(slide.get("title", "Slide"))
                text_frame = ppt_slide.placeholders[1].text_frame
                text_frame.clear()
                for index, bullet in enumerate(slide.get("bullets", [])):
                    paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                    paragraph.text = str(bullet)

            prs.save(str(output_path))
            return output_path
        except Exception:
            self._write_fallback_pptx(plan=plan, manuscript=manuscript, output_path=output_path)
            return output_path

    def _write_fallback_pptx(self, plan: dict[str, Any], manuscript: str, output_path: Path) -> None:
        title = _xml_escape(str(plan.get("title", "Presentation")))
        bullets = [
            _xml_escape(str(item))
            for slide in plan.get("slides", [])
            for item in ([slide.get("title", "")] + list(slide.get("bullets", [])))
            if str(item).strip()
        ]
        body_text = _xml_escape("\n".join(bullets) or manuscript or "Presentation generated by MyPaperWeb")
        slide_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr/>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Title"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
        <p:spPr/>
        <p:txBody>
          <a:bodyPr/>
          <a:lstStyle/>
          <a:p><a:r><a:t>{title}</a:t></a:r></a:p>
          <a:p><a:r><a:t>{body_text}</a:t></a:r></a:p>
        </p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>
"""
        slide_rels_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>
"""
        with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("[Content_Types].xml", _CONTENT_TYPES_XML)
            archive.writestr("_rels/.rels", _ROOT_RELS_XML)
            archive.writestr("ppt/presentation.xml", _PRESENTATION_XML)
            archive.writestr("ppt/_rels/presentation.xml.rels", _PRESENTATION_RELS_XML)
            archive.writestr("ppt/slides/slide1.xml", slide_xml)
            archive.writestr("ppt/slides/_rels/slide1.xml.rels", slide_rels_xml)
            archive.writestr("ppt/slideMasters/slideMaster1.xml", _SLIDE_MASTER_XML)
            archive.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", _SLIDE_MASTER_RELS_XML)
            archive.writestr("ppt/slideLayouts/slideLayout1.xml", _SLIDE_LAYOUT_XML)
            archive.writestr("ppt/theme/theme1.xml", _THEME_XML)
            archive.writestr("docProps/app.xml", _APP_XML)
            archive.writestr("docProps/core.xml", _CORE_XML.format(title=title))
